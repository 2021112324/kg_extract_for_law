import os
import logging
from typing import Dict, Any, List
from .base_processor import BaseProcessor, ProcessResult

logger = logging.getLogger(__name__)

class OfficeProcessor(BaseProcessor):
    """
    Office文档处理器
    
    处理Word、Excel、PowerPoint等Office文档
    支持.doc/.docx, .xls/.xlsx, .ppt/.pptx格式
    """
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
        
    async def process_file(self, file_path: str, output_dir: str, **kwargs) -> ProcessResult:
        """
        处理Office文档
        
        Args:
            file_path: Office文件路径
            output_dir: 输出目录
            **kwargs: 额外参数
            
        Returns:
            ProcessResult: 处理结果
        """
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            filename_uuid = kwargs.get('filename_uuid', 'office_output')
            
            # 提取文本内容
            extracted_text = await self.extract_text(file_path)
            
            # 保存提取的文本
            text_output_path = os.path.join(output_dir, f"{filename_uuid}_text.txt")
            os.makedirs(output_dir, exist_ok=True)
            
            with open(text_output_path, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return ProcessResult(
                success=True,
                file_path=file_path,
                file_type=file_ext,
                extracted_text=extracted_text,
                output_files=[text_output_path],
                metadata={
                    **self.get_metadata(file_path),
                    'document_type': self._get_document_type(file_ext),
                    'text_output_path': text_output_path,
                    'text_length': len(extracted_text)
                }
            )
            
        except Exception as e:
            error_msg = f"Office文档处理失败: {str(e)}"
            logger.error(error_msg)
            return self.create_error_result(file_path, error_msg)
    
    async def extract_text(self, file_path: str) -> str:
        """
        从 Office 文档中提取文本内容
        
        Args:
            file_path: Office文件路径
            
        Returns:
            str: 提取的文本内容
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext in ['.doc', '.docx']:
                return await self._extract_from_word(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                return await self._extract_from_excel(file_path)
            elif file_ext in ['.ppt', '.pptx']:
                return await self._extract_from_powerpoint(file_path)
            else:
                return f"不支持的Office文档类型: {file_ext}"
                
        except Exception as e:
            logger.error(f"Office文本提取失败: {str(e)}")
            return f"Office文本提取失败: {str(e)}"
    
    async def _extract_from_word(self, file_path: str) -> str:
        """从Word文档提取文本"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            return '\n'.join(text_content)
            
        except ImportError:
            logger.error("python-docx库未安装，无法处理Word文档")
            return "错误：python-docx库未安装"
        except Exception as e:
            logger.error(f"Word文档处理失败: {str(e)}")
            return f"Word文档处理失败: {str(e)}"
    
    async def _extract_from_excel(self, file_path: str) -> str:
        """从Excel文档提取文本"""
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_data = pd.read_excel(file_path, sheet_name=None)
            text_content = []
            
            for sheet_name, df in excel_data.items():
                text_content.append(f"=== 工作表: {sheet_name} ===")
                
                # 转换为文本格式
                if not df.empty:
                    # 添加列标题
                    headers = ' | '.join([str(col) for col in df.columns])
                    text_content.append(headers)
                    text_content.append('-' * len(headers))
                    
                    # 添加数据行
                    for _, row in df.iterrows():
                        row_text = ' | '.join([str(val) if pd.notna(val) else '' for val in row])
                        text_content.append(row_text)
                
                text_content.append('')  # 空行分隔
            
            return '\n'.join(text_content)
            
        except ImportError:
            logger.error("pandas库未安装，无法处理Excel文档")
            return "错误：pandas库未安装"
        except Exception as e:
            logger.error(f"Excel文档处理失败: {str(e)}")
            return f"Excel文档处理失败: {str(e)}"
    
    async def _extract_from_powerpoint(self, file_path: str) -> str:
        """从PowerPoint文档提取文本"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== 幻灯片 {i} ===")
                
                # 提取幻灯片中的文本
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                text_content.append('')  # 空行分隔
            
            return '\n'.join(text_content)
            
        except ImportError:
            logger.error("python-pptx库未安装，无法处理PowerPoint文档")
            return "错误：python-pptx库未安装"
        except Exception as e:
            logger.error(f"PowerPoint文档处理失败: {str(e)}")
            return f"PowerPoint文档处理失败: {str(e)}"
    
    def _get_document_type(self, file_ext: str) -> str:
        """获取文档类型描述"""
        type_mapping = {
            '.doc': 'Word文档 (DOC)',
            '.docx': 'Word文档 (DOCX)',
            '.xls': 'Excel表格 (XLS)',
            '.xlsx': 'Excel表格 (XLSX)',
            '.ppt': 'PowerPoint演示 (PPT)',
            '.pptx': 'PowerPoint演示 (PPTX)'
        }
        return type_mapping.get(file_ext.lower(), f'未知Office文档 ({file_ext})')
    
    async def extract_images_from_office(self, file_path: str, output_dir: str) -> List[str]:
        """
        从Office文档中提取图片（可选功能）
        
        Args:
            file_path: Office文件路径
            output_dir: 图片输出目录
            
        Returns:
            List[str]: 提取的图片文件路径列表
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        image_files = []
        
        try:
            os.makedirs(output_dir, exist_ok=True)
            
            if file_ext in ['.docx']:  # Word文档中的图片提取
                # 这里可以实现Word文档中图片的提取
                # 需要额外的库支持
                pass
            elif file_ext in ['.pptx']:  # PowerPoint中的图片提取
                # 这里可以实现PowerPoint中图片的提取
                pass
            
            logger.info(f"从Office文档提取了 {len(image_files)} 张图片")
            return image_files
            
        except Exception as e:
            logger.error(f"Office文档图片提取失败: {str(e)}")
            return []